"""
RAG (retrieval-augmented generation) service for the AI personal guide (Phase 6).

Runtime layer — requires Flask app context and the database. Orchestrates:
  - extracting text from CourseContent items (plain text, PDF, or video transcription)
  - chunking + embedding that text (via gemini_client) and storing it in ContentEmbedding
  - answering student questions by retrieving relevant chunks and grounding a Gemini
    response in them, with citations back to the source CourseContent

See gemini_client.py for the underlying Gemini REST calls (no Flask/DB dependency there,
mirroring core_translator.py's split from translation_service.py).
"""
import logging
import os
import re
import tempfile
from datetime import datetime, timedelta

import filetype
import requests
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from lms import gemini_client, r2_client, transcription
from lms.models import db, AiConversation, AiConversationMessage, ContentEmbedding, CourseContent

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1200
CHUNK_OVERLAP = 150

# Time-window target for video/audio chunking (chunk_segments_by_time) — independent of
# CHUNK_SIZE above, which only bounds chunk length for dense speech (see that function).
SEGMENT_CHUNK_WINDOW_SECONDS = 45.0

# Retrieval is two-stage: first rank whole files by their single best-matching chunk (from
# a wider candidate pool), then only pull context from the top few files. Mixing fragments
# from every vaguely-related file into one prompt gets noisy as a course accumulates more
# material — narrowing to the most relevant documents first keeps answers coherent.
CANDIDATE_POOL_SIZE = 30
MAX_SOURCE_FILES = 3
CHUNKS_PER_FILE = 4

# Conversation memory: keep this many raw turns (messages, not Q&A pairs) verbatim; older
# ones get folded into a rolling summary. A conversation the user hasn't consented to keep
# (AiConversation joined to a User with ai_history_consent True/False/None) is hard-deleted
# once inactive for this many days regardless of the above — see purge_stale_conversations.
MAX_VERBATIM_MESSAGES = 6
CONVERSATION_RETENTION_DAYS = 30

SYSTEM_INSTRUCTION = (
    "You are a study assistant answering questions about a specific course, using only the "
    "course material excerpts provided as context (plus prior conversation context, if any, "
    "purely to understand what the student is referring to). If the answer isn't contained "
    "in the provided excerpts, say you don't have that information in the course materials — "
    "do not use outside knowledge to fill gaps. Explain your answer thoroughly rather than "
    "giving a bare one-line response — walk through the relevant concepts, not just the "
    "conclusion. If the student's question is ambiguous, underspecified, or you'd need more "
    "detail from them to give a good answer, ask a concise clarifying follow-up question "
    "instead of guessing. Mention which material(s) you drew from by name. If an excerpt "
    "is labeled with a timestamp (e.g. 'at 5:12'), you may reference that moment in your "
    "answer (e.g. 'around the 5:12 mark') so the student can find it in the video."
)

# Effort modes: a student-facing toggle trading retrieval depth/answer length for speed. The
# instruction above (thorough explanations, follow-up questions) applies at both levels —
# this only controls how much source material gets pulled in and how much room the model has
# to write. `daily_cost` feeds a cost-weighted daily rate limit on /api/course/<id>/ask (see
# routes/api.py) — thorough genuinely costs more per call (more retrieval, thinking left on,
# 4x the output budget), and Gemini's actual free-tier daily cap on the underlying model was
# observed at 20 requests/day, so the shared daily budget is sized with that in mind rather
# than picked arbitrarily.
EFFORT_LEVELS = {
    'quick': {
        'max_files': MAX_SOURCE_FILES,
        'chunks_per_file': CHUNKS_PER_FILE,
        'max_output_tokens': 1024,
        'thinking_budget': 0,  # disabled - all of max_output_tokens goes to the visible answer
        'daily_cost': 1,
        'style_instruction': "Keep this particular answer relatively brief — a focused, direct explanation is enough here.",
    },
    'thorough': {
        'max_files': 5,
        'chunks_per_file': 6,
        'max_output_tokens': 4096,
        'thinking_budget': None,  # let the model think freely; budget above leaves it room to
        'daily_cost': 3,
        'style_instruction': (
            "Go deep on this one: cover the underlying concepts in detail, connect related "
            "points across the retrieved material, and use examples from the source content "
            "where they'd help understanding."
        ),
    },
}
DEFAULT_EFFORT = 'thorough'


# ── Chunking ────────────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Word-boundary-aware sliding-window chunking."""
    text = ' '.join(text.split())
    if not text:
        return []

    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        if end < n:
            break_point = text.rfind(' ', start, end)
            if break_point > start:
                end = break_point
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(end - overlap, start + 1)
    return chunks


def chunk_segments_by_time(
    segments: list[dict],
    window_seconds: float = SEGMENT_CHUNK_WINDOW_SECONDS,
    max_chars: int = CHUNK_SIZE,
) -> list[dict]:
    """Group timestamped transcript segments (from transcription.transcribe_with_timestamps)
    into chunks of ~window_seconds, each carrying the real start/end of the audio it covers.

    A window also closes early if it would exceed max_chars, so a chunk never grows unbounded
    over a long stretch of dense speech. Returns [{'text', 'start', 'end'}, ...].

    Degrades gracefully for the Gemini fallback path (_transcribe_via_gemini), which returns
    a single segment spanning the whole file with start=end=0.0 (no real timestamps available):
    the span-based close condition never fires since start never advances, but the max_chars
    condition still splits it into multiple reasonably-sized chunks — same (0.0, 0.0) on all
    of them, meaning "position unknown," rather than one unbounded blob.
    """
    chunks: list[dict] = []
    cur_texts: list[str] = []
    cur_start = None
    cur_end = None
    cur_chars = 0

    def flush():
        nonlocal cur_texts, cur_start, cur_end, cur_chars
        if cur_texts:
            chunks.append({'text': ' '.join(cur_texts), 'start': cur_start, 'end': cur_end})
        cur_texts, cur_start, cur_end, cur_chars = [], None, None, 0

    for seg in segments:
        text = (seg.get('text') or '').strip()
        if not text:
            continue

        if cur_texts and (
            seg['end'] - cur_start > window_seconds
            or cur_chars + len(text) + 1 > max_chars
        ):
            flush()

        if len(text) > max_chars:
            # A single segment's own text already exceeds max_chars — the Gemini-fallback
            # path hits this, returning one segment spanning the whole file. Flush whatever
            # was pending, then sub-split this segment's text with the same word-boundary
            # logic chunk_text() uses. All pieces share the segment's start/end, since a
            # single transcript segment carries no finer-grained timing to split by.
            flush()
            for piece in chunk_text(text, chunk_size=max_chars, overlap=0):
                chunks.append({'text': piece, 'start': seg['start'], 'end': seg['end']})
            continue

        if cur_start is None:
            cur_start = seg['start']
        cur_texts.append(text)
        cur_chars += len(text) + 1
        cur_end = seg['end']

    flush()
    return chunks



def _format_timestamp(seconds: float) -> str:
    """5:12, or 1:05:12 past the hour mark. Whole seconds only — sub-second precision isn't
    meaningful without a player that can actually seek to it."""
    total = int(seconds)
    h, rem = divmod(total, 3600)
    m, s = divmod(rem, 60)
    return f"{h}:{m:02d}:{s:02d}" if h else f"{m}:{s:02d}"

# ── Text extraction ──────────────────────────────────────────────────────────────

# Keyed by MIME, not extension — content titles are freeform display names (e.g. "Privacy
# policy" has no extension at all), so the actual file type has to come from sniffing the
# downloaded bytes (via `filetype`) rather than trusting the title. Found live: two real
# course files with extension-less titles were silently skipped entirely before any download
# was even attempted, because the old check looked at os.path.splitext(content.title).
_DOCUMENT_MIMES = {
    'application/pdf': '.pdf',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
}


def extract_text_for_content(content: CourseContent) -> str | list[dict] | None:
    """Best-effort content extraction for indexing. Returns None if the content type isn't
    supported (e.g. links, images, or legacy Office formats like .doc/.ppt — Office *preview
    rendering* is separately Phase 8 scope, but raw text extraction from the modern .docx/
    .pptx XML formats is unrelated to that and cheap enough to do here).

    Returns a plain str for untimed sources (text/PDF/DOCX/PPTX), or a list of timestamped
    segments (`[{'start','end','text'}, ...]`, see transcription.transcribe_with_timestamps)
    for video/audio — embed_content_item dispatches on which shape it got back."""
    if content.content_type == 'text':
        return content.content_data
    if content.content_type == 'file' and content.has_bytes:
        return _extract_drive_file_text(content)
    if content.content_type == 'video' and content.has_bytes:
        return transcribe_video_segments(content)
    return None


def _download_content_bytes(content: CourseContent, dest_path: str) -> bool:
    """Fetch a CourseContent's file bytes to dest_path from whichever backend holds them.
    R2 first (the source of truth once a row is migrated, per CourseContent.storage_backend),
    Drive as the fallback for rows not yet migrated off it."""
    if content.r2_key:
        return r2_client.download_file(content.r2_key, dest_path)
    if content.drive_file_id:
        return _download_public_drive_file(content.drive_file_id, dest_path)
    return False


def _download_public_drive_file(file_id: str, dest_path: str) -> bool:
    """Download a Drive file via its public share link, bypassing the OAuth-scoped Drive
    API entirely. Course files are made public-viewable at upload time (see
    google_drive_service.set_file_permissions) — the OAuth-scoped API (files.get/get_media)
    only works for files the querying identity itself created/opened (the drive.file scope's
    restriction), which 404s on anything created under a different identity than whichever
    the RAG pipeline authenticates as (e.g. content uploaded before the Phase 4 worker
    account existed). The public link has no such restriction, so this works uniformly."""
    session = requests.Session()
    url = 'https://drive.google.com/uc?export=download'
    try:
        resp = session.get(url, params={'id': file_id}, stream=True, timeout=60)

        # Large files get an HTML "can't scan for viruses" interstitial instead of the file
        # itself — needs a confirm token (from a cookie or the page body) to proceed.
        if 'text/html' in resp.headers.get('Content-Type', ''):
            token = next((v for k, v in resp.cookies.items() if k.startswith('download_warning')), None)
            if not token:
                match = re.search(r'confirm=([0-9A-Za-z_]+)', resp.text)
                token = match.group(1) if match else None
            if not token:
                logger.error(f"Could not extract confirm token for large Drive file {file_id}")
                return False
            resp = session.get(url, params={'id': file_id, 'confirm': token}, stream=True, timeout=300)

        if resp.status_code != 200:
            logger.error(f"Public Drive download for {file_id} returned {resp.status_code}")
            return False

        with open(dest_path, 'wb') as f:
            for chunk in resp.iter_content(chunk_size=1024 * 1024):
                if chunk:
                    f.write(chunk)
        return True
    except Exception as e:
        logger.error(f"Public Drive download failed for {file_id}: {e}")
        return False


def _extract_drive_file_text(content: CourseContent) -> str | list[dict] | None:
    """Returns extracted text for pdf/docx/pptx, or timestamped segments (list[dict]) for
    legacy rows mis-typed as 'file' that actually sniff as video/audio — see the self-healing
    branch below."""
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp_path = tmp.name
        if not _download_content_bytes(content, tmp_path):
            return None

        kind = filetype.guess(tmp_path)
        mime = kind.mime if kind else None

        # Self-healing path for legacy rows: content uploaded before content-type sniffing
        # existed is stored as content_type='file' even when it is really a lecture recording.
        # Rather than a backfill (which can't sniff Drive files from a migration), transcribe
        # it here off the copy we already downloaded — no second download, no manual fixup.
        if mime and (mime.startswith('video/') or mime.startswith('audio/')):
            logger.info(
                f"Content {content.id} is stored as 'file' but sniffs as {mime}; transcribing instead"
            )
            segments, language = transcription.transcribe_with_timestamps(tmp_path)
            content.transcript_language = language
            return segments

        extension = _DOCUMENT_MIMES.get(mime) if mime else None
        if not extension:
            return None

        if extension == '.pdf':
            reader = PdfReader(tmp_path)
            return '\n'.join(page.extract_text() or '' for page in reader.pages)
        if extension == '.docx':
            doc = DocxDocument(tmp_path)
            return '\n'.join(p.text for p in doc.paragraphs if p.text)
        if extension == '.pptx':
            prs = Presentation(tmp_path)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text:
                        lines.append(shape.text_frame.text)
            return '\n'.join(lines)
        return None
    except Exception as e:
        logger.error(f"Text extraction failed for content {content.id}: {e}")
        return None
    finally:
        if tmp_path:
            try:
                os.remove(tmp_path)
            except OSError:
                pass


def transcribe_video_segments(content: CourseContent) -> list[dict] | None:
    """Transcribe a video/audio CourseContent into timestamped segments.

    Returns [{'start': seconds, 'end': seconds, 'text': str}, ...] or None.
    Whisper is the primary engine (local, no quota, native timestamps); the older
    Gemini path stays as a fallback for when the model can't be loaded or fails.
    """
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp_path = tmp.name
        if not _download_content_bytes(content, tmp_path):
            return None

        # Keyed off the downloaded bytes, not content.title (see _extract_drive_file_text) —
        # a freeform display title ("Week 3 recording") can't be trusted to carry an extension.
        kind = filetype.guess(tmp_path)
        mime_type = kind.mime if kind else None
        if not mime_type or not (mime_type.startswith('video/') or mime_type.startswith('audio/')):
            return None

        segments, language = transcription.transcribe_with_timestamps(tmp_path)
        if segments:
            content.transcript_language = language
            return segments

        logger.warning(f"Whisper returned nothing for content {content.id}; trying Gemini fallback")
        text = _transcribe_via_gemini(tmp_path, mime_type, content.title)
        # Gemini gives no timestamps — represent it as a single unbounded segment so
        # callers get a consistent shape rather than a special case.
        return [{'start': 0.0, 'end': 0.0, 'text': text}] if text else None

    except Exception as e:
        logger.error(f"Video transcription failed for content {content.id}: {e}")
        return None
    finally:
        if tmp_path:
            try:
                os.remove(tmp_path)
            except OSError:
                pass


def _transcribe_via_gemini(tmp_path: str, mime_type: str, title: str) -> str | None:
    """Fallback transcription via Gemini's File API. No timestamps."""
    gemini_file_name = None
    try:
        file_uri, gemini_file_name = gemini_client.upload_file(tmp_path, mime_type, display_name=title)
        if not file_uri or not gemini_client.wait_for_file_active(gemini_file_name):
            return None
        return gemini_client.generate_content_with_file(
            file_uri, mime_type,
            'Transcribe this lecture recording verbatim. Output only the transcript text — '
            'no timestamps, no speaker labels, no commentary.',
            timeout=300,
        )
    finally:
        if gemini_file_name:
            gemini_client.delete_gemini_file(gemini_file_name)


# ── Indexing ────────────────────────────────────────────────────────────────────

def embed_content_item(content: CourseContent) -> int:
    """Extract, chunk, embed, and store one CourseContent item. Returns the number of
    chunks stored (0 if nothing was extracted, e.g. an unsupported file type).

    extract_text_for_content returns either a plain str (text/PDF/DOCX/PPTX — no time axis)
    or a list of timestamped segments (video/audio). Both are normalized here into the same
    {'text', 'start', 'end'} chunk shape before embedding, so the storage loop below doesn't
    need to know which source produced them — start/end are simply None for untimed chunks."""
    extracted = extract_text_for_content(content)

    ContentEmbedding.query.filter_by(course_content_id=content.id).delete()

    if isinstance(extracted, list):
        chunks = chunk_segments_by_time(extracted)
        # Stage 1 of video moment highlighting: an automatic keyword/regex pass over the raw
        # per-segment transcript (finer-grained than the 45s chunks above) — see
        # moment_service.py. Function-local import: moment_service imports from this module
        # at its own top level, so a module-level import here would cycle.
        from lms.moment_service import record_auto_moments
        record_auto_moments(content, extracted)
        # Subtitles: persist the same raw per-segment transcript at its original fine-grained
        # timing, separately from the 45s chunks above — see subtitle_service.py.
        from lms.subtitle_service import record_transcript_segments
        record_transcript_segments(content, extracted)
    elif extracted and extracted.strip():
        chunks = [{'text': t, 'start': None, 'end': None} for t in chunk_text(extracted)]
    else:
        chunks = []

    vectors = gemini_client.embed_batch([c['text'] for c in chunks], task_type='RETRIEVAL_DOCUMENT') if chunks else []

    stored = 0
    for idx, (chunk, vector) in enumerate(zip(chunks, vectors)):
        if not vector:
            continue
        db.session.add(ContentEmbedding(
            course_content_id=content.id,
            chunk_index=idx,
            chunk_text=chunk['text'],
            start_seconds=chunk['start'],
            end_seconds=chunk['end'],
            embedding=vector,
        ))
        stored += 1

    # Re-embed any already-captioned video moments (see moment_service.py) — no vision API
    # call, since the caption text is the durable artifact stored precisely so this is free.
    # Without this, the ContentEmbedding.query...delete() above would silently and
    # permanently drop every caption on the next reindex, since the promotion sweep never
    # re-captions a bucket that already has a VideoMoment row.
    from lms.moment_service import caption_chunks_for_content
    caption_chunks = caption_chunks_for_content(content)
    if caption_chunks:
        next_index = len(chunks)
        caption_vectors = gemini_client.embed_batch([c['text'] for c in caption_chunks], task_type='RETRIEVAL_DOCUMENT')
        for i, (chunk, vector) in enumerate(zip(caption_chunks, caption_vectors)):
            if not vector:
                continue
            db.session.add(ContentEmbedding(
                course_content_id=content.id,
                chunk_index=next_index + i,
                chunk_text=chunk['text'],
                start_seconds=chunk['start'],
                end_seconds=chunk['end'],
                embedding=vector,
            ))
            stored += 1

    content.embedded_at = datetime.utcnow()
    db.session.commit()
    return stored


# ── Retrieval + answering ────────────────────────────────────────────────────────

def get_locked_content_ids(course, user) -> set[int]:
    """Content ids inside a currently-locked folder for this user — mirrors the lock check
    in course_page_enrolled.html (folder.locked_until_assignment_id/locked_until_quiz_id vs.
    the user's passed submissions/attempts) so the assistant can't answer from gated
    material. Locks cascade to subfolders of a locked folder."""
    from lms.models import CourseAssignmentSubmission, Quiz, QuizAttempt

    if user is None or not getattr(user, 'is_authenticated', False):
        passed_assignment_ids = set()
        passed_quiz_ids = set()
    else:
        passed_assignment_ids = {
            s.assignment_id for s in
            CourseAssignmentSubmission.query.filter_by(user_id=user.id, passed=True).all()
        }
        passed_quiz_ids = {
            a.quiz_id for a in
            QuizAttempt.query.join(Quiz)
                .filter(QuizAttempt.user_id == user.id, Quiz.course_id == course.id, QuizAttempt.passed.is_(True))
                .all()
        }

    folders = {f.id: f for f in course.content_folders.all()}

    def folder_locked(folder_id, seen=frozenset()):
        if folder_id in seen or folder_id not in folders:
            return False
        folder = folders[folder_id]
        if folder.locked_until_assignment_id and folder.locked_until_assignment_id not in passed_assignment_ids:
            return True
        if folder.locked_until_quiz_id and folder.locked_until_quiz_id not in passed_quiz_ids:
            return True
        if folder.parent_folder_id:
            return folder_locked(folder.parent_folder_id, seen | {folder_id})
        return False

    return {
        content.id for content in course.contents.all()
        if content.folder_id and folder_locked(content.folder_id)
    }


def _get_or_create_conversation(user, course) -> AiConversation:
    conversation = AiConversation.query.filter_by(user_id=user.id, course_id=course.id).first()
    if not conversation:
        conversation = AiConversation(user_id=user.id, course_id=course.id)
        db.session.add(conversation)
        db.session.commit()
    return conversation


def _build_history_context(conversation: AiConversation) -> str:
    """Rolling summary (if any) + the most recent verbatim turns, formatted for inclusion
    in a prompt. Empty string if this is a fresh conversation."""
    parts = []
    if conversation.summary:
        parts.append(f"Summary of earlier conversation:\n{conversation.summary}")

    recent = (
        conversation.messages.order_by(AiConversationMessage.created_at.desc())
        .limit(MAX_VERBATIM_MESSAGES)
        .all()
    )
    recent.reverse()
    if recent:
        speaker = {'user': 'Student', 'assistant': 'Assistant'}
        turns = [f"{speaker[m.role]}: {m.content}" for m in recent]
        parts.append("Recent conversation:\n" + '\n'.join(turns))

    return '\n\n'.join(parts)


def _rewrite_query_for_retrieval(question: str, history_context: str) -> str:
    """Follow-up questions ('what about the second one?') often don't embed anywhere near
    the chunks they actually need — rewrite into a standalone question using conversation
    context before embedding for retrieval. Falls back to the original question on failure."""
    if not history_context:
        return question

    prompt = (
        f"Conversation so far:\n{history_context}\n\n"
        f"Latest student message: {question}\n\n"
        "Rewrite the latest student message as a standalone question that makes sense "
        "without the conversation history, resolving pronouns/references. If it's already "
        "standalone, return it unchanged. Output only the rewritten question, nothing else."
    )
    rewritten = gemini_client.generate_content(prompt, temperature=0.0, timeout=20)
    return rewritten.strip() if rewritten else question


def _persist_turn(conversation: AiConversation, question: str, answer: str, sources: list) -> None:
    db.session.add(AiConversationMessage(conversation_id=conversation.id, role='user', content=question))
    db.session.add(AiConversationMessage(
        conversation_id=conversation.id, role='assistant', content=answer, sources=sources or None,
    ))
    conversation.last_activity_at = datetime.utcnow()
    db.session.commit()
    _compact_conversation_if_needed(conversation)


def _compact_conversation_if_needed(conversation: AiConversation) -> None:
    """Fold the oldest turns into the rolling summary once the raw count grows past
    MAX_VERBATIM_MESSAGES, so a long-running conversation's prompt size stays bounded."""
    messages = conversation.messages.order_by(AiConversationMessage.created_at).all()
    if len(messages) <= MAX_VERBATIM_MESSAGES:
        return

    excess = messages[:len(messages) - MAX_VERBATIM_MESSAGES]
    speaker = {'user': 'Student', 'assistant': 'Assistant'}
    excess_text = '\n'.join(f"{speaker[m.role]}: {m.content}" for m in excess)

    prompt = (
        (f"Existing summary of the conversation so far:\n{conversation.summary}\n\n" if conversation.summary else "")
        + f"New exchanges to fold into the summary:\n{excess_text}\n\n"
        "Write an updated, concise summary (a few sentences) of the entire conversation so "
        "far, preserving important facts/topics discussed. Output only the summary text."
    )
    new_summary = gemini_client.generate_content(prompt, temperature=0.2, timeout=30)
    if not new_summary:
        return  # leave the raw messages in place and try again next turn

    conversation.summary = new_summary.strip()
    for m in excess:
        db.session.delete(m)
    db.session.commit()


def get_conversation_history(user, course) -> dict:
    """For the frontend on page load: {'consent': True/False/None, 'messages': [...]}.
    `messages` is only populated if the user has actually consented — the caller doesn't
    need to re-check consent itself, but `consent` is always returned so the frontend knows
    whether to show the opt-in prompt."""
    conversation = AiConversation.query.filter_by(user_id=user.id, course_id=course.id).first()
    messages = []
    if conversation and user.ai_history_consent:
        messages = [
            {'role': m.role, 'content': m.content, 'sources': m.sources or []}
            for m in conversation.messages.order_by(AiConversationMessage.created_at).all()
        ]
    return {'consent': user.ai_history_consent, 'messages': messages}


def purge_stale_conversations() -> int:
    """Hard-delete conversations (cascading their messages) the user hasn't consented to
    keep, once inactive for CONVERSATION_RETENTION_DAYS. Consented conversations are never
    auto-purged by this. Returns the number of conversations deleted."""
    from lms.models import User

    cutoff = datetime.utcnow() - timedelta(days=CONVERSATION_RETENTION_DAYS)
    stale = (
        AiConversation.query
        .join(User, AiConversation.user_id == User.id)
        .filter(
            AiConversation.last_activity_at < cutoff,
            db.or_(User.ai_history_consent.is_(None), User.ai_history_consent.is_(False)),
        )
        .all()
    )
    count = len(stale)
    for conversation in stale:
        db.session.delete(conversation)
    db.session.commit()
    return count


def answer_question(
    course, question: str, user,
    max_files: int | None = None, chunks_per_file: int | None = None,
    effort: str = DEFAULT_EFFORT,
) -> dict | None:
    """Returns {'answer': str, 'sources': [{'content_id', 'title', 'timestamp'?, 'start_seconds'?}]},
    or None if the question couldn't be processed (e.g. GEMINI_API_KEY missing/unreachable).
    Multiple entries may share a content_id — one per distinct moment actually cited in a
    video/audio file, distinguished by start_seconds — rather than one entry per file.
    Tracks multi-turn conversation memory per (user, course) when user is a real authenticated
    User — see _get_or_create_conversation / _persist_turn.

    `effort` ('quick'/'thorough', see EFFORT_LEVELS) sets retrieval depth and answer length
    defaults; explicit max_files/chunks_per_file, if passed, override the preset."""
    preset = EFFORT_LEVELS.get(effort, EFFORT_LEVELS[DEFAULT_EFFORT])
    if max_files is None:
        max_files = preset['max_files']
    if chunks_per_file is None:
        chunks_per_file = preset['chunks_per_file']

    conversation = None
    history_context = ''
    if user is not None and getattr(user, 'is_authenticated', False):
        conversation = _get_or_create_conversation(user, course)
        history_context = _build_history_context(conversation)

    retrieval_query = _rewrite_query_for_retrieval(question, history_context)

    query_vector = gemini_client.embed_text(retrieval_query, task_type='RETRIEVAL_QUERY')
    if not query_vector:
        return None

    locked_ids = get_locked_content_ids(course, user)

    query = (
        db.session.query(ContentEmbedding, CourseContent)
        .join(CourseContent, ContentEmbedding.course_content_id == CourseContent.id)
        .filter(
            CourseContent.course_id == course.id,
            CourseContent.is_published.is_(True),
        )
    )
    # Teacher-private files ("🔒 Private" in the content list) must not leak into answers.
    # Same rule the rest of the app applies — see the folder-contents endpoint and the
    # file-serving gate in routes/api.py: students only see allow_others_to_view=True,
    # course managers see everything. `is_managed_by` handles anonymous/None users.
    if not course.is_managed_by(user):
        query = query.filter(CourseContent.allow_others_to_view.is_(True))
    if locked_ids:
        query = query.filter(~CourseContent.id.in_(locked_ids))

    candidates = (
        query.order_by(ContentEmbedding.embedding.cosine_distance(query_vector))
        .limit(CANDIDATE_POOL_SIZE)
        .all()
    )

    if not candidates:
        answer = "I don't have any indexed course material to answer that from yet."
        sources = []
    else:
        # Stage 1: pick which files to read. candidates is already ordered by ascending
        # distance (best match first), so each file's first appearance here is its best
        # chunk — ranking files by that and keeping only the top few is the "file
        # selection" step.
        ranked_content_ids = []
        seen = set()
        for _, content in candidates:
            if content.id not in seen:
                seen.add(content.id)
                ranked_content_ids.append(content.id)
        selected_content_ids = set(ranked_content_ids[:max_files])

        # Stage 2: within just those files, take each file's best few chunks.
        context_blocks = []
        sources = []
        emitted_untimed = set()
        emitted_starts = {}  # content_id -> [start_seconds, ...] already given their own chip
        chunks_used = {}
        for embedding, content in candidates:
            if content.id not in selected_content_ids:
                continue
            if chunks_used.get(content.id, 0) >= chunks_per_file:
                continue
            chunks_used[content.id] = chunks_used.get(content.id, 0) + 1
            if embedding.start_seconds is not None:
                header = f"[Source: {content.title}, at {_format_timestamp(embedding.start_seconds)}]"
            else:
                header = f"[Source: {content.title}]"
            context_blocks.append(f"{header}\n{embedding.chunk_text}")

            # candidates is ordered by ascending distance, so the first chunk to claim a given
            # moment is that moment's best match. A video can earn several chips — one per
            # distinct moment actually cited — as long as each is more than one chunking
            # window away from every moment already emitted for that file; that keeps two
            # adjacent chunks of the same moment from spamming duplicate chips while still
            # surfacing genuinely separate moments the AI mentioned in the same answer.
            if embedding.start_seconds is None:
                if content.id not in emitted_untimed:
                    emitted_untimed.add(content.id)
                    sources.append({'content_id': content.id, 'title': content.title})
            else:
                starts = emitted_starts.setdefault(content.id, [])
                if all(abs(s - embedding.start_seconds) >= SEGMENT_CHUNK_WINDOW_SECONDS for s in starts):
                    starts.append(embedding.start_seconds)
                    sources.append({
                        'content_id': content.id,
                        'title': content.title,
                        'timestamp': _format_timestamp(embedding.start_seconds),
                        'start_seconds': embedding.start_seconds,
                    })

        # Files stay in relevance order; multiple moments within the same file read better
        # chronologically ("2:10, then 12:43") than in retrieval-score order.
        sources.sort(key=lambda s: (ranked_content_ids.index(s['content_id']), s.get('start_seconds', -1)))

        context = '\n\n---\n\n'.join(context_blocks)
        prompt_parts = []
        if history_context:
            prompt_parts.append(history_context)
        prompt_parts.append(f"Course material excerpts:\n\n{context}")
        prompt_parts.append(f"Student question: {question}")
        prompt = '\n\n---\n\n'.join(prompt_parts)

        system_instruction = f"{SYSTEM_INSTRUCTION}\n\n{preset['style_instruction']}"
        answer = gemini_client.generate_content(
            prompt,
            system_instruction=system_instruction,
            temperature=0.2,
            max_output_tokens=preset['max_output_tokens'],
            thinking_budget=preset['thinking_budget'],
        )
        if not answer:
            return None

    if conversation:
        _persist_turn(conversation, question, answer, sources)

    return {'answer': answer, 'sources': sources}
